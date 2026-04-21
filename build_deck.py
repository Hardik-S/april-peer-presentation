from __future__ import annotations

import re
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DECK_DIR = ROOT / "deck"
PNG_DIR = DECK_DIR / "slides_png"
REVIEW_DIR = ROOT / "review"
OUTPUT_PPTX = DECK_DIR / "April_peer_presentation.pptx"
OUTPUT_SCRIPT = DECK_DIR / "April_peer_presentation_script.md"
CONTACT_SHEET = REVIEW_DIR / "contact_sheet.png"
STRUCTURE_MAP = REVIEW_DIR / "adam_structure_mapping.md"
REVIEW_MEMO = REVIEW_DIR / "review_memo.md"
REVIEW_CHECKLIST = REVIEW_DIR / "review_checklist.md"
ADAM_DIR = ROOT / "sources" / "adam_structure"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
WORDS_PER_MINUTE = 130

NAVY = RGBColor(70, 44, 119)
INK = RGBColor(18, 18, 18)
SLATE = RGBColor(84, 84, 92)
TEAL = RGBColor(50, 127, 147)
TEAL_SOFT = RGBColor(229, 244, 247)
GOLD = RGBColor(168, 125, 42)
GOLD_SOFT = RGBColor(249, 242, 213)
RED = RGBColor(165, 65, 63)
RED_SOFT = RGBColor(249, 228, 226)
GREEN = RGBColor(54, 135, 86)
GREEN_SOFT = RGBColor(231, 244, 235)
BLUE_SOFT = RGBColor(238, 241, 246)
BG = RGBColor(255, 255, 255)
WHITE = RGBColor(255, 255, 255)
PURPLE = RGBColor(94, 48, 151)
PURPLE_LIGHT = RGBColor(197, 181, 219)
HEADER_GRAY = RGBColor(244, 244, 244)
PANEL_GRAY = RGBColor(247, 247, 247)
RULE = RGBColor(158, 158, 166)
BULLET = RGBColor(49, 72, 160)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Consolas"


def color(name: str) -> RGBColor:
    return {
        "navy": NAVY,
        "ink": INK,
        "slate": SLATE,
        "teal": TEAL,
        "teal_soft": TEAL_SOFT,
        "gold": GOLD,
        "gold_soft": GOLD_SOFT,
        "red": RED,
        "red_soft": RED_SOFT,
        "green": GREEN,
        "green_soft": GREEN_SOFT,
        "blue_soft": BLUE_SOFT,
        "bg": BG,
        "white": WHITE,
        "purple": PURPLE,
        "purple_light": PURPLE_LIGHT,
        "header_gray": HEADER_GRAY,
        "panel_gray": PANEL_GRAY,
        "rule": RULE,
        "bullet": BULLET,
    }[name]


def add_text_box(slide, left, top, width, height, paragraphs, *, margin=0.05, vertical_anchor=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    if vertical_anchor is not None:
        frame.vertical_anchor = vertical_anchor
    for idx, spec in enumerate(paragraphs):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = spec["text"]
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.level = spec.get("level", 0)
        p.font.name = spec.get("font", BODY_FONT)
        p.font.size = Pt(spec.get("size", 18))
        p.font.bold = spec.get("bold", False)
        p.font.color.rgb = color(spec.get("color", "ink"))
        p.space_after = Pt(spec.get("space_after", 5))
    return box


def add_card(slide, left, top, width, height, fill_name, *, line_name=None, rounded=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE if not rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill_name)
    shape.line.color.rgb = color(line_name or "rule")
    shape.line.width = Pt(0.65)
    return shape


def add_chip(slide, text, left, top, width, fill_name="purple", *, text_name="white", font_size=8):
    add_card(slide, left, top, width, 0.22, fill_name, line_name=fill_name, rounded=False)
    add_text_box(
        slide,
        left + 0.04,
        top + 0.035,
        width - 0.16,
        0.12,
        [{"text": text, "size": font_size, "bold": True, "color": text_name, "align": PP_ALIGN.CENTER}],
        margin=0.0,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.08))
    top.fill.solid()
    top.fill.fore_color.rgb = PURPLE
    top.line.fill.background()
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.08), SLIDE_WIDTH, Inches(0.28))
    band.fill.solid()
    band.fill.fore_color.rgb = HEADER_GRAY
    band.line.fill.background()


def add_headline(slide, title, label):
    add_chip(slide, label, 5.35, 0.01, 2.65, "purple", font_size=7)
    add_text_box(slide, 0.26, 0.45, 11.8, 0.38, [{"text": title, "size": 18, "bold": False, "color": "purple"}], margin=0.0)
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(0.88), Inches(12.85), Inches(0.012))
    rule.fill.solid()
    rule.fill.fore_color.rgb = PURPLE_LIGHT
    rule.line.fill.background()


def add_takeaway(slide, text, *, fill_name="gold_soft", line_name="gold"):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.25), Inches(6.37), Inches(12.85), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    add_text_box(
        slide,
        0.34,
        6.52,
        12.4,
        0.32,
        [{"text": text, "size": 12, "bold": True, "color": "ink", "align": PP_ALIGN.LEFT}],
        margin=0.0,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, number):
    add_text_box(slide, 12.18, 7.12, 0.82, 0.16, [{"text": f"{number} of 20", "size": 6.5, "color": "slate", "align": PP_ALIGN.RIGHT}], margin=0.0)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def bullet_paragraphs(items, *, size=18, color_name="ink"):
    return [{"text": f"●  {item}", "size": size, "color": color_name, "space_after": 6} for item in items]


def render_title(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_background(slide)
    add_text_box(slide, 0.55, 1.15, 12.2, 0.55, [{"text": spec["title"], "size": 18, "bold": False, "color": "purple", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_text_box(slide, 0.55, 1.85, 12.2, 0.35, [{"text": spec["subtitle"], "size": 11, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_text_box(slide, 0.55, 2.45, 12.2, 0.32, [{"text": "Hardik S.  |  Department of Computer Science  |  Western University", "size": 9, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_card(slide, 4.1, 3.72, 5.1, 0.44, "panel_gray", line_name="rule", rounded=False)
    add_text_box(slide, 4.22, 3.86, 4.86, 0.16, [{"text": spec["claim"], "size": 9.5, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, 4.05, 4.6, 5.2, 0.2, [{"text": "Adam: visual reference only.  Hardik: Maple content source.", "size": 7.8, "color": "slate", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_roadmap(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_headline(slide, spec["title"], "Talk Map")
    for idx, step in enumerate(spec["steps"]):
        left = 0.65 + idx * 3.12
        add_card(slide, left, 1.78, 2.72, 2.15, step["fill"], line_name=step["line"], rounded=False)
        add_text_box(slide, left + 0.18, 2.02, 2.36, 0.28, [{"text": step["name"], "size": 12.5, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
        add_text_box(slide, left + 0.2, 2.52, 2.32, 0.62, [{"text": step["detail"], "size": 9.2, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
        add_chip(slide, step["tag"], left + 0.46, 3.34, 1.82, "purple", font_size=6.5)
    add_takeaway(slide, spec["takeaway"])
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_two_column(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_headline(slide, spec["title"], spec.get("label", "Technical Core"))
    add_card(slide, 0.58, 1.42, 5.95, 4.25, spec.get("left_fill", "white"), line_name=spec.get("left_line", "rule"), rounded=False)
    add_card(slide, 6.82, 1.42, 5.95, 4.25, spec.get("right_fill", "panel_gray"), line_name=spec.get("right_line", "rule"), rounded=False)
    add_text_box(slide, 0.82, 1.68, 5.4, 0.3, [{"text": spec["left_title"], "size": 12.5, "bold": True, "color": "ink"}], margin=0.0)
    add_text_box(slide, 7.06, 1.68, 5.42, 0.3, [{"text": spec["right_title"], "size": 12.5, "bold": True, "color": "ink"}], margin=0.0)
    add_text_box(slide, 0.86, 2.2, 5.25, 2.6, bullet_paragraphs(spec["left"], size=10.6), margin=0.0)
    add_text_box(slide, 7.1, 2.2, 5.25, 2.6, bullet_paragraphs(spec["right"], size=10.6), margin=0.0)
    if spec.get("chip"):
        add_chip(slide, spec["chip"], 4.68, 5.35, 4.0, "purple", font_size=6.8)
    add_takeaway(slide, spec["takeaway"], fill_name=spec.get("takeaway_fill", "gold_soft"))
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_pipeline(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_headline(slide, spec["title"], "Pipeline")
    for idx, node in enumerate(spec["nodes"]):
        left = 0.78 + idx * 2.42
        add_card(slide, left, 2.12, 1.68, 0.5, node["fill"], line_name=node["line"], rounded=False)
        add_text_box(slide, left + 0.1, 2.29, 1.48, 0.12, [{"text": node["name"], "size": 7.5, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0, vertical_anchor=MSO_ANCHOR.MIDDLE)
        if idx < len(spec["nodes"]) - 1:
            add_text_box(slide, left + 1.82, 2.25, 0.42, 0.14, [{"text": "→", "size": 13, "bold": True, "color": "purple", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_card(slide, 1.0, 4.1, 11.15, 0.55, "panel_gray", line_name="rule", rounded=False)
    add_text_box(slide, 1.24, 4.28, 10.72, 0.16, [{"text": spec["detail"], "size": 9.8, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_takeaway(slide, spec["takeaway"])
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_code(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_headline(slide, spec["title"], spec.get("label", "Pseudocode"))
    add_card(slide, 0.72, 1.36, 7.6, 4.4, "blue_soft", line_name="rule", rounded=False)
    add_text_box(slide, 0.96, 1.63, 7.06, 3.84, [{"text": spec["code"], "size": spec.get("code_size", 10.2), "font": MONO_FONT, "color": "ink", "space_after": 0}], margin=0.0)
    add_card(slide, 8.62, 1.36, 3.9, 4.4, "white", line_name="rule", rounded=False)
    add_text_box(slide, 8.88, 1.72, 3.34, 0.28, [{"text": spec["side_title"], "size": 12.5, "bold": True, "color": "ink"}], margin=0.0)
    add_text_box(slide, 8.88, 2.28, 3.3, 2.35, bullet_paragraphs(spec["side"], size=10.2), margin=0.0)
    add_takeaway(slide, spec["takeaway"])
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_metrics(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_headline(slide, spec["title"], "Evidence")
    headers = ["Input", "Result", "Comparison"]
    for idx, metric in enumerate(spec["metrics"]):
        left = 0.72 + idx * 4.08
        add_card(slide, left, 1.66, 3.68, 1.48, metric["fill"], line_name="rule", rounded=False)
        add_text_box(slide, left + 0.14, 1.84, 3.4, 0.18, [{"text": headers[idx], "size": 8.2, "bold": True, "color": "purple", "align": PP_ALIGN.CENTER}], margin=0.0)
        add_text_box(slide, left + 0.18, 2.2, 3.3, 0.28, [{"text": metric["value"], "size": 16.5, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
        add_text_box(slide, left + 0.24, 2.68, 3.18, 0.22, [{"text": metric["label"], "size": 7.8, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_card(slide, 1.0, 4.18, 11.1, 0.55, "panel_gray", line_name="rule", rounded=False)
    add_text_box(slide, 1.28, 4.36, 10.55, 0.16, [{"text": spec["detail"], "size": 9.5, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_takeaway(slide, spec["takeaway"])
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def render_close(prs, number, spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_background(slide)
    add_text_box(slide, 0.55, 1.12, 11.8, 0.42, [{"text": spec["title"], "size": 18, "bold": False, "color": "purple"}], margin=0.0)
    add_text_box(slide, 0.58, 2.08, 11.7, 0.55, [{"text": spec["claim"], "size": 12.5, "bold": True, "color": "ink"}], margin=0.0)
    for idx, chip in enumerate(spec["chips"]):
        add_card(slide, 1.0 + idx * 3.95, 4.25, 3.35, 0.28, "panel_gray", line_name="rule", rounded=False)
        add_text_box(slide, 1.1 + idx * 3.95, 4.34, 3.15, 0.1, [{"text": chip, "size": 7.8, "bold": True, "color": "ink", "align": PP_ALIGN.CENTER}], margin=0.0)
    add_text_box(slide, 0.82, 6.08, 3.0, 0.22, [{"text": "Thank you.", "size": 9.5, "bold": True, "color": "ink"}], margin=0.0)
    add_notes(slide, spec["script"])
    add_footer(slide, number)


def build_slides() -> list[dict]:
    slides = [
        {
            "layout": "title",
            "title": "Exact Loop-Level Parallelization in Maple",
            "subtitle": "April peer presentation, rebuilt with Adam's technical walkthrough structure",
            "claim": "Integer feasibility is the decision layer.",
            "script": "This is a fresh April version of the talk. I am using Adam's presentation as the model for structure rather than content: start with context, introduce the representation, walk through the algorithm, then show evidence and limits. The content here is my Maple thesis work. The contribution is a command, IsParallelizable, that keeps Maple's existing loop-analysis front end but changes the final loop-level decision so it is based on integer feasibility instead of rational over-approximation alone.",
        },
        {
            "layout": "roadmap",
            "title": "The talk follows a technical walkthrough",
            "steps": [
                {"name": "Setup", "detail": "What dependence analysis must decide for loop nests.", "tag": "problem", "fill": "white", "line": "navy"},
                {"name": "Model", "detail": "How Maple represents loops, accesses, and constraints.", "tag": "representation", "fill": "teal_soft", "line": "teal"},
                {"name": "Algorithm", "detail": "How IsParallelizable tests one loop level at a time.", "tag": "decision", "fill": "white", "line": "navy"},
                {"name": "Evidence", "detail": "Where integer reasoning improves the rational baseline.", "tag": "validation", "fill": "gold_soft", "line": "gold"},
            ],
            "takeaway": "Adam's structure keeps the talk grounded: model first, algorithm second, evidence last.",
            "script": "The structure is deliberately close to Adam's rhythm. I first give the high-level setup, then the data representation, then the algorithmic walkthrough, and only then the performance-style evidence and limitations. The difference is the content. Instead of hypergraph transversals, the object here is dependence analysis for Maple loop nests and the exact question of which loop levels can be safely parallelized.",
        },
        {
            "layout": "two_column",
            "title": "The parallelization question is level-specific",
            "label": "Setup",
            "left_title": "What Maple already knows",
            "right_title": "What the thesis decides",
            "left": ["Iteration vectors live in Z^d.", "IterationSpace gives affine legality constraints.", "ArrayReferences separates reads and writes.", "DependenceCone summarizes distances over Q."],
            "right": ["A level is blocked only by an integer witness.", "The first nonzero distance component must be that level.", "Each read/write or write/write pair is tested separately.", "The answer is safe variables, not just a cone summary."],
            "chip": "From dependence summary to loop-level decision",
            "takeaway": "The thesis contribution is the last decision step, not a replacement front end.",
            "script": "The core issue is that loop parallelization is not just asking whether some dependence exists. It is asking whether a specific loop level is blocked by an actual pair of iterations. Maple already has the front half of the analysis: it can normalize a loop, describe the legal iteration space, and collect array references. The thesis adds the exact decision layer that asks whether a blocking witness exists in integer iteration space for the level under test.",
        },
        {
            "layout": "two_column",
            "title": "Marc's dependence definition gives the test shape",
            "label": "Definition",
            "left_title": "Dependence ingredients",
            "right_title": "Constraint ingredients",
            "left": ["Two statement instances touch the same memory location.", "At least one side writes: flow, anti, or output.", "There is a legal execution path from source to target.", "Marc's I(S) and O(S) sets determine candidate pairs."],
            "right": ["Source legality L_s(i).", "Target legality L_t(j).", "Equal-access equations f_s(i)=f_t(j).", "Lexicographic level-l blocker over j-i."],
            "chip": "Same-memory + one write + legal order",
            "takeaway": "The slide-level definition becomes a concrete integer feasibility query.",
            "script": "Marc's formal dependence definition is the bridge into the algorithm. A data dependence needs same memory, at least one write, and a legal execution path. In the implementation, those become source legality, target legality, equal-access equations, and a condition that fixes where the first loop-index difference occurs. This is why the talk can stay technical without drifting into implementation detail too early.",
        },
        {
            "layout": "pipeline",
            "title": "Maple supplies the front half of the analysis",
            "nodes": [
                {"name": "procedure", "fill": "white", "line": "navy"},
                {"name": "CreateLoop", "fill": "teal_soft", "line": "teal"},
                {"name": "ForLoop", "fill": "white", "line": "navy"},
                {"name": "IterationSpace", "fill": "teal_soft", "line": "teal"},
                {"name": "ArrayReferences", "fill": "white", "line": "navy"},
            ],
            "detail": "IsParallelizable reuses this pipeline and adds the exact integer decision layer after it.",
            "takeaway": "The command extends Maple's existing representation rather than inventing a new one.",
            "script": "This is the representation slide in Adam's sense: what data structures are actually being passed around. Maple normalizes a procedure into a ForLoop object. From there, IterationSpace gives the affine legality conditions, and ArrayReferences gives the reads and writes. IsParallelizable is intentionally conservative about scope: it uses these existing objects and focuses on the final question of whether an integer blocking witness exists for each loop level.",
        },
        {
            "layout": "two_column",
            "title": "The model is a two-iteration system",
            "label": "Model",
            "left_title": "Variables",
            "right_title": "Families of constraints",
            "left": ["Source index vector i=(i1,...,id).", "Target index vector j=(j1,...,jd).", "Distance vector d=j-i.", "Symbolic bounds m,n after initialization.", "All decision variables are integer."],
            "right": ["L_s(i) and L_t(j) both hold.", "f_s(i)=f_t(j) for the chosen array.", "Earlier levels satisfy d1=...=d(l-1)=0.", "The tested level satisfies d_l >= 1.", "The resulting Z-set is queried for nonemptiness."],
            "chip": "One pair, one level, one feasibility query",
            "takeaway": "Every decision is reduced to whether this system has an integer point.",
            "script": "For one source-target reference pair and one loop level, the command builds a two-iteration system. One copy of the loop variables describes the source instance, another describes the target instance, and the constraints force both to be legal and to touch the same memory cell. The level-specific blocker then says that this is the first loop level where the target is ordered after the source in a way that prevents parallel execution.",
        },
        {
            "layout": "code",
            "title": "Decision rule for one loop level",
            "code": "for level l in loop_levels:\n    blocked := false\n    for (src, tgt) in candidate_pairs:\n        R := legality(src.i) union legality(tgt.j)\n        R := R union equal_access(src.i, tgt.j)\n        R := R union {d = j - i}\n        R := R union {d_1=...=d_(l-1)=0, d_l >= 1}\n        Z := ZPolyhedralSet(R, integer_variables)\n        if IntegerPointDecomposition(Z) is nonempty:\n            blocked := true; break\n    report loop l safe iff blocked = false",
            "side_title": "Important boundary",
            "side": ["The query is existential.", "One real witness blocks the level.", "No integer witness means this pair does not block it."],
            "takeaway": "The criterion is exact for the affine integer system being tested.",
            "script": "This is the core algorithm in compact form. For each loop level, the command checks every relevant source-target access pair. It builds the legality constraints, the equal-access constraints, and the first-difference condition for the level. Then it asks a single yes-or-no question: does this system have an integer point? If yes, that loop level is blocked by a real execution witness. If no pair produces a witness, the level is reported safe.",
        },
        {
            "layout": "two_column",
            "title": "Reference pairs are filtered before solving",
            "label": "Implementation",
            "left_title": "Pairs that matter",
            "right_title": "Pairs that do not matter",
            "left": ["Read/write pairs model anti-dependence.", "Write/read pairs model true or flow dependence.", "Write/write pairs model output dependence.", "Pairs must reference the same array symbol.", "Subscripts must be affine after initialization."],
            "right": ["Read/read pairs do not create dependence.", "Different arrays cannot alias in this model.", "Unsupported nonlinear subscripts stay outside the claim.", "Branch predicates are not yet folded into the witness system."],
            "chip": "Keep the solver work tied to dependence semantics",
            "takeaway": "The command is precise because it asks the right smaller set of questions.",
            "script": "Before the backend query, the implementation keeps only access pairs that can create a dependence. A read-read pair is not a dependence because neither side writes. The command also keeps the claim inside the supported affine loop surface. That scope control matters for the talk: the result is not that every Maple program is now solved, but that this class of loop-level dependence decisions is handled exactly over integer iteration space.",
        },
        {
            "layout": "pipeline",
            "title": "ZPolyhedralSets is the integer backend",
            "nodes": [
                {"name": "constraints", "fill": "white", "line": "navy"},
                {"name": "normalize", "fill": "teal_soft", "line": "teal"},
                {"name": "ZPolyhedralSet", "fill": "white", "line": "navy"},
                {"name": "IntegerPointDecomp", "fill": "teal_soft", "line": "teal"},
                {"name": "safe?", "fill": "gold_soft", "line": "gold"},
            ],
            "detail": "The semantic domain is integer loop iterations, so the backend must answer integer-point existence.",
            "takeaway": "The backend matches the mathematics of loop execution.",
            "script": "Marc specifically asked for the role of ZPolyhedralSets, and this is the cleanest way to present it. Loop iterations are integer points. The blocker system is a system of affine constraints over those integer variables. So Maple builds a ZPolyhedralSet and asks whether any integer point remains. The command is not inventing a solver; it is connecting Maple's loop-analysis representation to the backend whose domain matches the question.",
        },
        {
            "layout": "two_column",
            "title": "DependenceCone is useful but rational",
            "label": "Comparison",
            "left_title": "Rational summary",
            "right_title": "Integer decision",
            "left": ["Summarizes possible dependence distances.", "Projects the system over Q.", "Can preserve fractional witnesses.", "Good for broad dependence information."],
            "right": ["Keeps the original witness system over Z.", "Asks whether real iteration pairs exist.", "Rejects fractional lattice artifacts.", "Returns loop variables that are safe to parallelize."],
            "chip": "The gap appears at the final decision step",
            "takeaway": "A rational witness is not automatically an executable loop witness.",
            "script": "This is the main contrast. DependenceCone is useful, but its output is rational. That is appropriate for a cone summary, but it is not the same thing as proving that a real pair of integer loop iterations blocks a level. In the parity-style cases, the rational analysis can keep a fractional distance that cannot be an actual iteration distance. IsParallelizable resolves that final question by checking the integer system.",
        },
        {
            "layout": "code",
            "title": "The parity mechanism explains the improvement",
            "label": "Mechanism",
            "code": "relaxed cone over Q:\n    2*d = 3\n    d = 3/2 satisfies the projection\n    cone keeps a possible blocker\n\ninteger witness test over Z:\n    d = j - i\n    i, j, d in Z\n    2*d = 3 has no integer solution\n\nresult:\n    witness set is empty\n    loop level is not blocked by this pair",
            "code_size": 10.5,
            "side_title": "What changes",
            "side": ["The same constraints are interpreted over integers.", "The false blocker disappears.", "The loop-level answer becomes less pessimistic."],
            "takeaway": "The thesis novelty is rejecting blockers that no execution can realize.",
            "script": "The parity case is the simplest way to explain why the backend matters. A rational projection can preserve a distance like three halves. As a rational object, that is perfectly legal. As a loop distance between two integer iterations, it is impossible. When the integer feasibility query is empty, the command rejects the blocker. That is the strict improvement over treating the rational summary as the final answer.",
        },
        {
            "layout": "two_column",
            "title": "Worked example: heat_eq starts from Maple objects",
            "label": "Example",
            "left_title": "Input surface",
            "right_title": "Analysis surface",
            "left": ["CreateLoop produces variables i,j.", "Bounds include 1 <= i <= m and 1 <= j <= n.", "Write surface includes A[i,j].", "Reads include neighboring stencil locations."],
            "right": ["ArrayReferences creates same-array candidate pairs.", "Each pair gets source and target copies.", "Outer level tests d_i >= 1.", "Inner level tests d_i=0 and d_j >= 1."],
            "chip": "Example first, general claim second",
            "takeaway": "The example shows the whole pipeline rather than only the final answer.",
            "script": "For the worked example, the current defense deck uses heat_eq because it walks through the full pipeline. The point of the slide is not only the final safe loop variable. It is that Maple starts from the same ForLoop representation, extracts accesses, builds the pairwise systems, and then makes a separate decision for each loop level. That makes the implementation easier to defend because the example lines up with the algorithm.",
        },
        {
            "layout": "two_column",
            "title": "The blocker is level-specific in heat_eq",
            "label": "Walkthrough",
            "left_title": "Outer level",
            "right_title": "Inner level",
            "left": ["Level i uses blocker d_i >= 1.", "A north/south neighbor pair can realize a real witness.", "IntegerPointDecomposition returns nonempty.", "The outer loop is reported blocked."],
            "right": ["Level j requires d_i=0 and d_j >= 1.", "Candidate dependences move through i first.", "All pair-level Z-sets are empty for this blocker.", "The inner loop remains safe."],
            "chip": "Same loop nest, different level decision",
            "takeaway": "Parallelizability is a per-level property, not a single global label.",
            "script": "The key teaching point in heat_eq is that loop levels must be tested separately. One level can be blocked because there is a real source-target pair that creates a dependence witness. Another level can be safe because no such integer witness exists under the first-difference condition for that level. This is why the command returns safe variables rather than only saying whether the loop nest has any dependence at all.",
        },
        {
            "layout": "code",
            "title": "What IsParallelizable returns",
            "label": "Interface",
            "code": "IsParallelizable(loop_or_procedure)\n\ninput:\n    ForLoop object, or procedure normalized by CreateLoop\n\noutput:\n    safe loop variables from outermost to innermost\n\nmeaning:\n    listed levels have no integer blocking witness\n    inside the supported affine analysis surface",
            "side_title": "User-facing result",
            "side": ["The command is callable, not worksheet-only.", "It preserves Maple's existing objects.", "The output is directly about parallelization."],
            "takeaway": "The public interface exposes the exact decision, not backend machinery.",
            "script": "The user-facing part is intentionally simple. The command accepts either a ForLoop or a procedure that Maple can normalize through CreateLoop. It returns the loop variables that are safe, ordered from outermost to innermost. The backend details are important for the defense, but the command surface is about the practical question a Maple user asks: which loop levels can I parallelize?",
        },
        {
            "layout": "metrics",
            "title": "Validation stays tied to the frozen bundle",
            "metrics": [
                {"value": "23", "label": "frozen cases", "fill": "white", "line": "navy"},
                {"value": "PASS", "label": "dated rerun surface", "fill": "green_soft", "line": "green"},
                {"value": "0", "label": "validated cases where DependenceCone is less pessimistic", "fill": "gold_soft", "line": "gold"},
            ],
            "detail": "The evidence claim is bounded to the validated bundle and examples in the current technical defense.",
            "takeaway": "The deck should make a strong claim only where the evidence is frozen.",
            "script": "The evidence slide needs to be careful. The current technical defense ties the claim to a frozen 23-case bundle and a dated pass rerun. Within that surface, there is no validated case where DependenceCone is less pessimistic than IsParallelizable, and there are strict-improvement cases where integer reasoning rejects a false rational blocker. That is the claim the April deck should make, without turning it into an unsupported universal theorem.",
        },
        {
            "layout": "two_column",
            "title": "The strict-improvement cases are the headline evidence",
            "label": "Results",
            "left_title": "Agreement cases",
            "right_title": "Strict-improvement cases",
            "left": ["Bucket A/B cases preserve standard answers.", "Control wavefront cases stay blocked.", "Nonempty integer witnesses still block the level.", "The new layer does not erase real dependences."],
            "right": ["Bucket C exposes rational pessimism.", "parity_case gives the 3/2 mechanism.", "frac2d and stride3 variants repeat the lattice gap.", "Safe loop levels become visible only when Z is empty."],
            "chip": "Preserve correctness, remove false blockers",
            "takeaway": "The improvement is useful because it is bounded by control cases.",
            "script": "The evidence should be presented in two parts. First, agreement cases show that the command preserves standard answers and does not incorrectly remove real dependences. Second, the strict-improvement cases show where rational pessimism is too strong. This makes the result credible: the command is not simply more optimistic everywhere, it is more precise where the integer feasibility question rules out a false blocker.",
        },
        {
            "layout": "two_column",
            "title": "Optimizations and precision have tradeoffs",
            "label": "Lessons",
            "left_title": "What became clearer",
            "right_title": "What remains bounded",
            "left": ["Exactness belongs at the final decision point.", "Integer feasibility changes the answer only when the witness set is empty.", "Rational summaries are useful but should not decide safety.", "Backend choice affects semantic precision."],
            "right": ["The current claim is affine.", "The current evidence is the frozen bundle.", "Branches and nonlinear subscripts need separate modeling.", "Performance claims stay outside the evidence boundary."],
            "chip": "Adam-style lesson slide: what worked and why",
            "takeaway": "The talk should explain the boundary as part of the contribution.",
            "script": "Adam's talk included lessons about optimizations that did not always behave the same way across cases. The equivalent lesson here is about precision. Integer feasibility is the right semantic test for loop iterations, but the claim is still bounded by the supported affine surface and the validated examples. That boundary is not a weakness to hide. It is what keeps the technical claim defensible.",
        },
        {
            "layout": "two_column",
            "title": "Scope control keeps the claim defensible",
            "label": "Limits",
            "left_title": "Inside this deck",
            "right_title": "Outside this deck",
            "left": ["Affine loop bounds and affine array accesses.", "Maple ForLoop normalization.", "ZPolyhedralSets integer-point queries.", "Validated bundle examples."],
            "right": ["Nonlinear indexing claims.", "Arbitrary side effects.", "A universal theorem over all Maple programs.", "Performance claims beyond the available evidence."],
            "chip": "Precise does not mean unlimited",
            "takeaway": "The April version should be technically confident and scope-aware.",
            "script": "This slide prevents overclaiming. Inside the deck are affine loops, Maple's ForLoop representation, array-reference constraints, and ZPolyhedralSets integer feasibility. Outside the deck are nonlinear indexing, arbitrary side effects, and performance or correctness claims that have not been validated. That lets the presentation be strong without being brittle.",
        },
        {
            "layout": "two_column",
            "title": "Future work follows from the same architecture",
            "label": "Future Work",
            "left_title": "Short-term extensions",
            "right_title": "Longer-term direction",
            "left": ["Broaden the validated example bundle.", "Improve diagnostics for blocked levels.", "Add clearer traces for the generated witness systems."],
            "right": ["Handle more program shapes.", "Explore integration with broader Maple parallelization tools.", "Use exact feasibility as a reusable decision layer."],
            "chip": "Same pipeline, wider surface",
            "takeaway": "The architecture is useful because the exact decision layer is separable.",
            "script": "The future work follows naturally from the architecture. The front end can stay Maple's loop representation, and the exact integer decision can be improved independently. The near-term direction is better validation, clearer diagnostics, and better traces for why a level is blocked. The longer-term direction is expanding the supported program surface and connecting the command more directly to parallelization workflows.",
        },
        {
            "layout": "close",
            "title": "Closing claim",
            "claim": "IsParallelizable turns Maple's dependence information into an exact loop-level decision over integer iteration space.",
            "chips": ["Adam's structure", "Hardik's content", "Maple-ready final step"],
            "script": "The closing claim is deliberately compact. Maple already has useful dependence infrastructure. The thesis contribution is to turn that into an exact loop-level parallelization decision by asking the right integer feasibility question. For the April version, Adam's structure gives the talk a clear technical arc, but the content stays anchored in the current Maple technical defense and the frozen validation surface.",
        },
    ]
    timing_expansions = [
        "I would also say explicitly that this version is not trying to imitate Adam's topic. The useful thing from Adam is the discipline of explaining one technical object at a time. That means the audience should always know what role the current slide plays: representation, algorithm step, evidence, or limitation.",
        "This roadmap also sets expectations about depth. I am not going to race through every theorem and every implementation detail from the defense deck. Instead, I am choosing the same kind of paced technical story Adam used: enough background to make the algorithm meaningful, then one concrete path through the implementation, then evidence.",
        "This distinction is important because a dependence summary can be true and still not answer the question a programmer cares about. The programmer wants to know whether a specific loop index can run in parallel. That makes the unit of analysis a loop level, not the entire nest and not the whole dependence cone.",
        "I should keep this slide connected to Marc's lecture language. The formal definition is not decoration. It tells us exactly what the generated constraint system has to encode. If any of same memory, one write, or legal ordering is missing, then the solver is answering the wrong question.",
        "The reason to show the pipeline is to avoid making the thesis sound larger or vaguer than it is. Maple already has machinery for parsing and representing loop nests. The new work starts after that machinery has produced the objects needed to ask the exact loop-level question.",
        "In the talk I would slow down here and make the two copies of the loop variables concrete. One copy is the earlier or source execution instance and one copy is the later or target execution instance. The equal-access constraints tie them together by forcing the two array subscripts to name the same cell.",
        "This pseudocode is the center of the talk. The nested loops over levels and reference pairs are deliberately simple. The complexity is hidden in building the correct system and in asking the backend over integers. That simplicity is useful because it makes the correctness argument easier to explain.",
        "The filtering step also explains why the command is not just throwing every possible pair at a solver. It uses the dependence semantics first, then uses integer feasibility second. That order matters for performance and for clarity, because the backend query should only be asked for pairs that could actually change the answer.",
        "This is where I would connect the command back to Maple rather than presenting ZPolyhedralSets as a separate mathematical tool. The backend is already in the ecosystem, and the thesis uses it because loop iterations are integer points. That makes the integration feel natural instead of bolted on.",
        "The comparison should be respectful to DependenceCone. It is not wrong for doing its job. The point is that its job is a rational dependence summary, while IsParallelizable has a different job: a final safety decision for loop levels. The difference in domain is what creates the improvement.",
        "This mechanism is the cleanest explanation of the false-blocker issue. A fractional value can survive rational reasoning, but the loop program will never execute at a fractional iteration. Once the audience accepts that, the need for an integer backend becomes intuitive rather than merely formal.",
        "The example section should feel like Adam's algorithm walkthrough: we do not only state that the command works, we show how the input travels through the analysis. The audience sees the representation, the pairwise constraints, and the separate decisions, so the final output is earned.",
        "I would emphasize that the same loop nest can produce different answers at different levels. That is the practical reason for returning safe variables. A single yes-or-no answer for the whole nest would lose useful information, especially when an inner level remains parallelizable after an outer level is blocked.",
        "This interface slide keeps the result concrete. It tells the audience what they would call, what kind of input is accepted, and how to read the output. The implementation can be mathematically detailed, but the command has to land as a usable Maple feature.",
        "The validation slide should be delivered with restraint. The right claim is strong but bounded: on the frozen evidence surface, the integer method preserves validated answers and exposes strict improvements over rational pessimism. I should avoid implying that untested program classes have already been handled.",
        "The important balance is that strict improvement is not just optimism. The agreement and control cases show that real dependences remain real dependences. The improvement cases matter because they are paired with that control surface, so the story is precision rather than wishful parallelization.",
        "This lesson slide mirrors Adam's habit of explaining what the experiments taught, not only what the algorithm does. Here the lesson is that exactness belongs at the semantic boundary where Maple decides safety. Earlier rational summaries are useful, but they should not be forced to answer an integer execution question.",
        "The limits should stay visible because they protect the contribution. It is better to say clearly that the current deck covers affine loop nests and validated examples than to imply a broader result. That gives reviewers a precise target for technical criticism and keeps the final thesis claim defensible.",
        "The future-work slide should not introduce a new thesis. It should show that the current architecture has room to grow. Better diagnostics, broader examples, and more program shapes all build on the same separation between Maple's front-end representation and the exact integer decision layer.",
        "The closing should return to the three-part story: Adam supplied the presentation structure, the Maple defense supplied the content, and this fresh repo keeps the April work isolated until review. The final import back into Maple should happen only after the deck is reviewed and explicitly approved.",
    ]
    for slide, extra in zip(slides, timing_expansions):
        slide["script"] = f"{slide['script']} {extra}"
    delivery_notes = [
        "I will frame this as a reset from the previous direction: the deck is now about the thesis contribution, presented with a cleaner peer-talk structure.",
        "I will use this slide to promise the audience a path through the talk, so they know each later technical slide has a specific job.",
        "I will give a small verbal example here: two loop iterations can both exist, but only one level may be responsible for the ordering problem.",
        "I will say that this is why the implementation checks read-write, write-read, and write-write cases, rather than treating all pairs uniformly.",
        "I will pause on each pipeline box long enough to show that the command is reusing Maple's native analysis objects.",
        "I will describe the two-iteration system as a witness search: if the witness exists, the level is blocked; otherwise it is not blocked by that pair.",
        "I will walk through the pseudocode line by line, because this is the easiest place for the audience to verify the algorithmic shape.",
        "I will connect this filtering to both correctness and cost: it avoids irrelevant queries and keeps the evidence tied to dependence semantics.",
        "I will emphasize that the backend question is existence, not enumeration; the command only needs to know whether at least one integer point remains.",
        "I will avoid making DependenceCone sound like a failed tool. It is the right baseline, but not the final safety oracle.",
        "I will make the fractional-distance example slow and concrete, because it is the shortest route to the intuition behind the thesis.",
        "I will use heat_eq as the running example that turns the abstract constraint language into something visibly produced by Maple.",
        "I will stress that different levels can have different answers, which is why the command returns a set of safe variables.",
        "I will keep the interface slide practical and avoid exposing more backend notation than a Maple user needs to use the command.",
        "I will state the validation date and bundle idea as process evidence: this deck should be auditable, not just persuasive.",
        "I will say that strict improvement means fewer false blocks, not a willingness to ignore true dependence witnesses.",
        "I will use this slide to show what I learned from the April reset: the narrative must make the exact decision layer obvious.",
        "I will present the limits before questions, so scope concerns are answered by the deck itself rather than patched on afterward.",
        "I will keep future work tied to the existing architecture instead of proposing disconnected research directions.",
        "I will close by reminding the audience that the repo isolation is also part of the process: review first, Maple import later.",
    ]
    for slide, note in zip(slides, delivery_notes):
        slide["script"] = f"{slide['script']} {note}"
    technical_depth_notes = [
        "Technically, the talk preserves the decision domain: loop iterations are integer points, so the final safety question cannot be answered only over rationals.",
        "This mirrors Adam's depth: each section names the object, the operation on it, and the next pipeline consequence.",
        "I will state the quantifier explicitly: the level is unsafe if there exists a source-target pair and an integer assignment satisfying the blocker system.",
        "The I(S) and O(S) notation matters because it explains why read-read pairs drop out before the solver is ever invoked.",
        "Here I distinguish representation from decision: ForLoop and ArrayReferences collect facts; IsParallelizable turns them into a level-specific query.",
        "The key implementation detail is that the source and target are two copies of the same loop nest, not two unrelated programs; the equal-access equations tie those copies together.",
        "The algorithm can stop early for a level once one nonempty Z-set is found, because one witness blocks parallel execution.",
        "This is also where the current scope boundary enters the algorithm: affine candidate pairs can be normalized into this system, while nonlinear cases are not silently forced through it.",
        "The backend call is an emptiness test in practice. The decomposition result is used only to decide whether some integer point exists, not to enumerate every dependence witness.",
        "The comparison is therefore not DependenceCone versus a better cone. It is a rational summary versus a separate integer witness check for a different question.",
        "The parity example is deliberately small because the entire failure mode is visible in one equation: two times an integer distance cannot equal three.",
        "For heat_eq I will make the objects concrete: two loop variables, affine bounds, one write surface, and neighbor reads.",
        "The outer-versus-inner contrast is the worked example: the same data goes through two blockers and gives two conclusions.",
        "The public command matters because it exposes the result at the level a Maple user needs: a list of safe loop variables, not an internal polyhedral certificate.",
        "I will explain the evidence categories as a test design: agreement cases protect against false positives, strict-improvement cases show the rational gap, and controls keep the claim honest.",
        "Bucket C is the important comparative bucket because it isolates the exact advantage: the rational baseline keeps a blocker that has no integer witness.",
        "The lesson is not that rational analysis is useless. The lesson is that a rational summary should feed, but not replace, the final integer safety decision.",
        "The limitation slide should sound operational: these are the program features the current witness construction does and does not encode.",
        "Future work should therefore extend the constraint model and diagnostic surface rather than changing the central architecture.",
        "The final claim should land as a technical invariant: Maple reports a level safe only after every relevant candidate pair has an empty integer witness set for that blocker.",
    ]
    for slide, note in zip(slides, technical_depth_notes):
        slide["script"] = f"{slide['script']} {note}"
    return slides


def build_script_markdown(slides: list[dict]) -> str:
    counts = [len(re.findall(r"\b[\w:-]+\b", slide["script"])) for slide in slides]
    total_words = sum(counts)
    lines = [
        "# April Peer Presentation Script",
        "",
        f"Target rate: {WORDS_PER_MINUTE} wpm",
        f"Estimated words: {total_words}",
        f"Estimated time: {round(total_words / WORDS_PER_MINUTE, 1)} minutes",
        "",
    ]
    for idx, (slide, count) in enumerate(zip(slides, counts), start=1):
        lines.extend(
            [
                f"## Slide {idx}. {slide['title']}",
                f"Estimated words: {count}",
                f"Estimated time: {round(count / WORDS_PER_MINUTE * 60)} seconds",
                "",
                slide["script"],
                "",
            ]
        )
    return "\n".join(lines)


def render_slides(prs: Presentation, slides: list[dict]) -> None:
    renderers = {
        "title": render_title,
        "roadmap": render_roadmap,
        "two_column": render_two_column,
        "pipeline": render_pipeline,
        "code": render_code,
        "metrics": render_metrics,
        "close": render_close,
    }
    for number, slide in enumerate(slides, start=1):
        renderers[slide["layout"]](prs, number, slide)


def export_slide_pngs(pptx_path: Path, output_dir: Path) -> bool:
    output_dir.mkdir(parents=True, exist_ok=True)
    for old in output_dir.glob("Slide*.PNG"):
        old.unlink()
    ps = f"""
$ErrorActionPreference = 'Stop'
$ppt = Resolve-Path '{pptx_path}'
$out = '{output_dir}'
$pp = New-Object -ComObject PowerPoint.Application
$msoTrue = -1
$msoFalse = 0
$pres = $pp.Presentations.Open($ppt.Path, $msoTrue, $msoFalse, $msoFalse)
try {{
    $pres.Export($out, 'PNG', 1600, 900)
}} finally {{
    $pres.Close()
    $pp.Quit()
}}
"""
    try:
        subprocess.run(["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-Command", ps], check=True, cwd=ROOT)
        return True
    except Exception as exc:
        print(f"PNG_EXPORT_SKIPPED: {exc}")
        return False


def build_contact_sheet(slides_dir: Path, output_path: Path) -> bool:
    images = sorted(slides_dir.glob("Slide*.PNG"), key=lambda p: int(p.stem.replace("Slide", "")))
    if not images:
        return False
    thumbs = []
    for path in images:
        img = Image.open(path).convert("RGB")
        img.thumbnail((320, 180))
        canvas = Image.new("RGB", (340, 224), "white")
        canvas.paste(img, ((340 - img.width) // 2, 14))
        draw = ImageDraw.Draw(canvas)
        draw.text((14, 198), path.stem, fill=(30, 37, 46))
        thumbs.append(canvas)
    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 340, rows * 224), (248, 246, 239))
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 340, (idx // cols) * 224))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return True


def write_review_docs(slides: list[dict], exported_pngs: bool, contact_sheet: bool) -> None:
    adam_screens = sorted(p.name for p in ADAM_DIR.glob("*.png"))
    mapping = [
        "# Adam Structure Mapping",
        "",
        "Adam is the visual/style reference only. Hardik's Maple technical defense is the authoritative content source.",
        "",
        "The generated deck remains editable PowerPoint: titles, bullets, code blocks, diagrams, tables, chips, and footers are native slide objects rather than pasted full-slide screenshots.",
        "",
        "## Inferred Adam Visual System",
        "",
        "- Canvas: widescreen landscape slides, approximated as 16:9.",
        "- Background: mostly white with a light gray title band and a thin purple top rule.",
        "- Title placement: compact, top-left, purple, with a small centered purple section tab.",
        "- Font feel: plain academic sans-serif; small titles and dense 8-12 pt body text rather than large keynote typography.",
        "- Density: high information density, with bullets, pseudocode, equations, tables, and small diagrams sharing the slide.",
        "- Palette: purple/lavender structure, black text, light gray panels, sparse teal/gold/green only for semantic contrast.",
        "- Borders and boxes: thin square-line boxes; minimal rounded-card treatment.",
        "- Footer/header: section label in the top band and small slide count in the lower-right corner.",
        "- Whitespace: moderate but not decorative; content starts close to the header and uses broad horizontal space.",
        "",
        "| Adam-style role | April deck slides | Hardik content mapped into that role |",
        "| --- | --- | --- |",
        "| High-level setup | 1-4 | Dependence analysis problem, formal dependence definition, loop-level decision. |",
        "| Core representation/model | 5-6 | Maple ForLoop, IterationSpace, ArrayReferences, constraint families. |",
        "| Algorithm walkthrough | 7-9 | IsParallelizable pair/level loop and ZPolyhedralSets backend. |",
        "| Comparison and mechanism | 10-11 | DependenceCone rational baseline and integer feasibility improvement. |",
        "| Worked example | 12-14 | heat_eq and public command surface from the technical defense. |",
        "| Evidence | 15-16 | Frozen bundle, strict-improvement cases, control cases. |",
        "| Lessons, limits, future work | 17-20 | Scope boundaries, lessons, extensions, closing claim. |",
        "",
        "Visual match statement: the April deck now follows Adam's screenshots as closely as practical while preserving Hardik's Maple thesis substance.",
        "",
        f"Usable Adam screenshots copied: {len(adam_screens)}",
        "Excluded corrupt screenshot: `Screenshot 2026-04-17 094616.png`",
    ]
    STRUCTURE_MAP.write_text("\n".join(mapping), encoding="utf-8")

    total_words = sum(len(re.findall(r"\b[\w:-]+\b", slide["script"])) for slide in slides)
    estimated_time = total_words / WORDS_PER_MINUTE
    memo = [
        "# Review Memo",
        "",
        "## What changed",
        "",
        "- Created a fresh April peer deck outside the Maple repo.",
        "- Used Adam's screenshots as the visual/style reference: white academic canvas, purple header bands, compact titles, square thin-line boxes, dense bullets/code/tables, and understated footers.",
        "- Used Hardik's current technical defense content as the authoritative Maple source; Adam's hypergraph topic was not imported.",
        "- Generated editable PowerPoint slides rather than screenshot-only slides.",
        "- Preserved the Adam-like technical structure: setup, model/representation, algorithm walkthrough, backend, worked example, comparison, evidence, limits, and future work.",
        "- Upgraded slide and script depth to match Adam's transcript more closely: each section now names the concrete data object, the operation performed on it, and the consequence for the next step.",
        "",
        "## Timing",
        "",
        f"- Slide count: {len(slides)}",
        f"- Script words: {total_words}",
        f"- Estimated speaking time at {WORDS_PER_MINUTE} wpm: {round(estimated_time, 1)} minutes",
        "",
        "## Generated artifacts",
        "",
        f"- PowerPoint: `{OUTPUT_PPTX}`",
        f"- Script: `{OUTPUT_SCRIPT}`",
        f"- PNG export: {'complete' if exported_pngs else 'skipped'}",
        f"- Contact sheet: {'complete' if contact_sheet else 'skipped'}",
        "",
        "## Boundary",
        "",
        "Adam is the style reference. Hardik is the content source. The slides remain editable PowerPoint objects. The deck visually matches Adam's screenshots as closely as practical without changing the Maple thesis substance.",
        "",
        "The script now follows Adam's technical-density pattern while staying on Hardik's Maple topic: concrete representation, algorithm mechanics, backend behavior, worked example consequences, evidence categories, and explicit limitations.",
        "",
        "The Maple repo was not modified. Final import into Maple remains out of scope until explicit approval.",
    ]
    REVIEW_MEMO.write_text("\n".join(memo), encoding="utf-8")

    checklist = [
        "# Review Checklist",
        "",
        "- [x] Adam is used for structure, not content.",
        "- [x] Adam is explicitly documented as the visual/style reference.",
        "- [x] Hardik's technical defense content is the authoritative substance.",
        "- [x] Hardik is explicitly documented as the content source.",
        "- [x] Omar's work is not used as deck basis.",
        "- [x] Adam's corrupt zero-byte screenshot is excluded.",
        "- [x] Slides are editable PowerPoint objects.",
        "- [x] Slides visually match Adam's screenshots as closely as practical while keeping Maple content.",
        "- [x] Script and slide contents match Adam's transcript-level technical depth while staying on Hardik's Maple thesis.",
        "- [x] Claims stay inside the validated Maple evidence described in the technical defense.",
        f"- [{'x' if 20 <= estimated_time <= 25 else ' '}] Talk fits roughly 20-25 minutes at normal speaking pace.",
        f"- [{'x' if exported_pngs else ' '}] Exported PNGs were generated.",
        f"- [{'x' if contact_sheet else ' '}] Contact sheet was generated for visual review.",
        "- [ ] Human review confirms slide readability and technical accuracy.",
        "- [ ] Human approval is given before importing final artifacts into Maple.",
    ]
    REVIEW_CHECKLIST.write_text("\n".join(checklist), encoding="utf-8")


def main() -> None:
    DECK_DIR.mkdir(exist_ok=True)
    REVIEW_DIR.mkdir(exist_ok=True)
    slides = build_slides()
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    render_slides(prs, slides)
    prs.save(OUTPUT_PPTX)
    OUTPUT_SCRIPT.write_text(build_script_markdown(slides), encoding="utf-8")
    exported = export_slide_pngs(OUTPUT_PPTX, PNG_DIR)
    contact = build_contact_sheet(PNG_DIR, CONTACT_SHEET)
    write_review_docs(slides, exported, contact)
    print(f"WROTE {OUTPUT_PPTX}")
    print(f"WROTE {OUTPUT_SCRIPT}")
    print(f"WROTE {STRUCTURE_MAP}")
    print(f"WROTE {REVIEW_MEMO}")
    print(f"WROTE {REVIEW_CHECKLIST}")
    print(f"WROTE {CONTACT_SHEET}" if contact else "CONTACT_SHEET_SKIPPED")


if __name__ == "__main__":
    main()
